from typing import Optional, Dict
import PyPDF2
from docx import Document
from pptx import Presentation
from io import BytesIO

class DocumentProcessor:
    def __init__(self):
        pass
    
    def process_document(self, file_content: bytes, filename: str) -> Dict:
        """Process uploaded document and extract text"""
        file_ext = filename.lower().split('.')[-1]
        
        if file_ext == 'pdf':
            text = self._extract_from_pdf(file_content)
        elif file_ext in ['docx', 'doc']:
            text = self._extract_from_docx(file_content)
        elif file_ext in ['pptx', 'ppt']:
            text = self._extract_from_pptx(file_content)
        elif file_ext == 'txt':
            text = file_content.decode('utf-8')
        else:
            text = file_content.decode('utf-8', errors='ignore')
        
        return {
            "filename": filename,
            "file_type": file_ext,
            "text": text,
            "word_count": len(text.split()),
            "char_count": len(text)
        }
    
    def _extract_from_pdf(self, content: bytes) -> str:
        """Extract text from PDF"""
        pdf_reader = PyPDF2.PdfReader(BytesIO(content))
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() or ""
        return text
    
    def _extract_from_docx(self, content: bytes) -> str:
        """Extract text from DOCX"""
        doc = Document(BytesIO(content))
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text
    
    def _extract_from_pptx(self, content: bytes) -> str:
        """Extract text from PPTX"""
        prs = Presentation(BytesIO(content))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    
    def extract_sections(self, text: str) -> list:
        """Extract sections from document text"""
        sections = []
        lines = text.split('\n')
        current_section = {"title": "Introduction", "content": ""}
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            if line.isupper() or (len(line) < 100 and line.endswith(':')):
                if current_section["content"]:
                    sections.append(current_section)
                current_section = {"title": line, "content": ""}
            else:
                current_section["content"] += line + " "
        
        if current_section["content"]:
            sections.append(current_section)
        
        return sections if sections else [{"title": "Content", "content": text}]
